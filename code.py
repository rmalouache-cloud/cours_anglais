import streamlit as st
import json
import os
from pathlib import Path
import time
from PIL import Image, ImageDraw, ImageFont
import io
import base64
import subprocess
import shutil

# ==============================
# PAGE CONFIG
# ==============================
st.set_page_config(
    page_title="✨ English Teacher's Platform ✨",
    page_icon="🌸",
    layout="wide"
)

# ==============================
# CUSTOM CSS
# ==============================
st.markdown("""
<style>
    .stApp {
        background: linear-gradient(135deg, #ffe6f0 0%, #ffd9e8 100%);
    }
    h1, h2, h3 {
        color: #c2185b !important;
    }
    .stButton > button {
        background: linear-gradient(45deg, #ff69b4, #ff1493);
        color: white;
        border-radius: 25px;
        border: none;
        padding: 12px 25px;
        font-weight: bold;
        transition: all 0.3s ease;
    }
    .stButton > button:hover {
        transform: scale(1.05);
    }
    .course-card {
        background: white;
        border-radius: 20px;
        padding: 20px;
        margin: 10px 0;
        box-shadow: 0 5px 15px rgba(0,0,0,0.08);
        border: 1px solid #ffc0cb;
    }
    @keyframes fadeInUp {
        from { transform: translateY(20px); opacity: 0; }
        to   { transform: translateY(0);    opacity: 1; }
    }
    .fade-in { animation: fadeInUp 0.6s ease-out; }

    /* Navigation bar for slides */
    .slide-nav {
        display: flex;
        align-items: center;
        justify-content: center;
        gap: 16px;
        margin: 12px 0;
    }
    .slide-counter {
        background: white;
        border: 2px solid #ff69b4;
        border-radius: 20px;
        padding: 6px 20px;
        color: #c2185b;
        font-weight: bold;
        font-size: 16px;
        min-width: 120px;
        text-align: center;
    }
    /* Progress bar */
    .prog-wrap {
        background: #ffe0ec;
        border-radius: 10px;
        height: 10px;
        margin: 4px 0 14px 0;
        overflow: hidden;
    }
    .prog-fill {
        height: 100%;
        border-radius: 10px;
        background: linear-gradient(90deg, #ff69b4, #ff1493);
        transition: width .3s ease;
    }
    /* Slide image wrapper */
    .slide-frame {
        background: white;
        border-radius: 18px;
        padding: 12px;
        box-shadow: 0 8px 30px rgba(255,20,147,.15);
        border: 2px solid #ffc0cb;
        text-align: center;
    }
    .slide-frame img { border-radius: 10px; max-width: 100%; }
</style>
""", unsafe_allow_html=True)

# ==============================
# FOLDERS INIT
# ==============================
def init_folders():
    for level in ["A", "B", "C"]:
        for sub in ["1", "2", "3"]:
            Path(f"courses/Level_{level}/{level}{sub}").mkdir(parents=True, exist_ok=True)
    Path("data").mkdir(exist_ok=True)

# ==============================
# METADATA
# ==============================
def load_metadata():
    path = "data/courses_metadata.json"
    if os.path.exists(path):
        with open(path, "r") as f:
            return json.load(f)
    return {}

def save_metadata(metadata):
    with open("data/courses_metadata.json", "w") as f:
        json.dump(metadata, f, indent=4)

# ==============================
# DELETE COURSE
# ==============================
def delete_course(course_key, course_path):
    try:
        if os.path.exists(course_path):
            os.remove(course_path)
        # Remove cached PNG folder
        cache_dir = Path(course_path).parent / (Path(course_path).stem + "_slides")
        if cache_dir.exists():
            shutil.rmtree(cache_dir)
        metadata = load_metadata()
        if course_key in metadata:
            del metadata[course_key]
            save_metadata(metadata)
        return True
    except Exception as e:
        st.error(f"Delete error: {e}")
        return False

# ==============================
# PPT → LIST OF PIL IMAGES
# Strategy priority:
#   1. LibreOffice  → PDF → pdf2image (best fidelity)
#   2. python-pptx  → manual render (fallback, no extra binaries)
# ==============================

def _libreoffice_to_images(ppt_path: Path) -> list[Image.Image] | None:
    """Convert via LibreOffice + pdf2image. Best quality."""
    # Find LibreOffice
    lo = None
    for candidate in ["libreoffice", "soffice",
                       "/usr/bin/libreoffice", "/usr/bin/soffice",
                       "/Applications/LibreOffice.app/Contents/MacOS/soffice"]:
        if shutil.which(candidate):
            lo = candidate
            break
    if not lo:
        return None

    pdf_path = ppt_path.with_suffix(".pdf")
    try:
        subprocess.run(
            [lo, "--headless", "--convert-to", "pdf",
             "--outdir", str(ppt_path.parent), str(ppt_path)],
            capture_output=True, timeout=120, check=True
        )
    except Exception:
        return None

    if not pdf_path.exists():
        return None

    # Try pdf2image
    try:
        from pdf2image import convert_from_path
        images = convert_from_path(str(pdf_path), dpi=150)
        return images
    except Exception:
        return None


def _pptx_to_images(ppt_path: Path) -> list[Image.Image]:
    """
    Pure python-pptx + Pillow renderer.
    Renders backgrounds, embedded images AND text with correct positions.
    Works without LibreOffice or poppler.
    """
    from pptx import Presentation

    DPI        = 150
    EMU        = 914400  # EMUs per inch

    prs = Presentation(str(ppt_path))
    slide_w = int(prs.slide_width  / EMU * DPI)
    slide_h = int(prs.slide_height / EMU * DPI)

    slides_out = []

    for slide in prs.slides:
        canvas = Image.new("RGB", (slide_w, slide_h), "white")

        # ── Background colour ──
        try:
            bg   = slide.background
            fill = bg.fill
            if str(fill.type) == "SOLID":
                rgb = fill.fore_color.rgb
                canvas = Image.new("RGB", (slide_w, slide_h),
                                   (rgb[0], rgb[1], rgb[2]))
        except Exception:
            pass

        draw = ImageDraw.Draw(canvas)

        for shape in slide.shapes:
            try:
                lft = int(shape.left  / EMU * DPI)
                top = int(shape.top   / EMU * DPI)
                w   = int(shape.width / EMU * DPI)
                h   = int(shape.height/ EMU * DPI)

                # ── Embedded image ──
                if hasattr(shape, "image"):
                    img = Image.open(io.BytesIO(shape.image.blob)).convert("RGBA")
                    if w > 0 and h > 0:
                        img = img.resize((w, h), Image.LANCZOS)
                        canvas.paste(img, (lft, top),
                                     img if img.mode == "RGBA" else None)

                # ── Text ──
                elif hasattr(shape, "text_frame"):
                    tf = shape.text_frame
                    y  = top

                    for para in tf.paragraphs:
                        line_parts = []
                        font_size  = 18
                        color      = (0, 0, 0)
                        bold       = False

                        for run in para.runs:
                            line_parts.append(run.text)
                            try:
                                if run.font.size:
                                    font_size = max(10, int(run.font.size.pt * DPI / 72))
                                if run.font.bold:
                                    bold = True
                                if run.font.color and run.font.color.type is not None:
                                    rgb   = run.font.color.rgb
                                    color = (rgb[0], rgb[1], rgb[2])
                            except Exception:
                                pass

                        text = "".join(line_parts).strip()
                        if not text:
                            y += font_size + 4
                            continue

                        # Try to load a real font
                        font = None
                        font_paths = [
                            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else
                            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                            "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
                            "/System/Library/Fonts/Helvetica.ttc",
                        ]
                        for fp in font_paths:
                            try:
                                font = ImageFont.truetype(fp, font_size)
                                break
                            except Exception:
                                pass
                        if font is None:
                            font = ImageFont.load_default()

                        # Word-wrap within shape width
                        words   = text.split()
                        lines   = []
                        cur     = []
                        for word in words:
                            test = " ".join(cur + [word])
                            bbox = draw.textbbox((0, 0), test, font=font)
                            if bbox[2] - bbox[0] <= w or not cur:
                                cur.append(word)
                            else:
                                lines.append(" ".join(cur))
                                cur = [word]
                        if cur:
                            lines.append(" ".join(cur))

                        lh = font_size + 5
                        for line in lines:
                            if y + lh > top + h + 20:
                                break
                            draw.text((lft, y), line, font=font, fill=color)
                            y += lh

            except Exception:
                continue

        slides_out.append(canvas)

    return slides_out


def get_slide_images(ppt_path: str) -> list[Image.Image]:
    """
    Return cached list of PIL Images for a PPT file.
    Cache is stored as PNGs next to the source file.
    """
    ppt  = Path(ppt_path)
    cache_dir = ppt.parent / (ppt.stem + "_slides")

    # Check cache freshness
    if cache_dir.exists():
        pngs = sorted(cache_dir.glob("slide_*.png"))
        if pngs and cache_dir.stat().st_mtime >= ppt.stat().st_mtime:
            return [Image.open(p) for p in pngs]

    # Regenerate
    cache_dir.mkdir(parents=True, exist_ok=True)

    # Try best method first
    images = _libreoffice_to_images(ppt)
    if not images:
        images = _pptx_to_images(ppt)

    # Save cache
    for i, img in enumerate(images):
        img.save(str(cache_dir / f"slide_{i:04d}.png"), format="PNG")

    return images


# ==============================
# DISPLAY PRESENTATION  (FIXED)
# ==============================
def display_presentation(course):
    st.markdown('<div class="fade-in">', unsafe_allow_html=True)

    if st.button("◀ Back to Courses"):
        st.session_state["viewing_course"] = None
        for k in ["slide_images", "current_course_path", "slide_index"]:
            st.session_state.pop(k, None)
        st.rerun()

    st.markdown(f"### 📖 {course['title']}")
    st.markdown(f"🎯 Level **{course['level']}** &nbsp;|&nbsp; 📅 {course['upload_date']}")
    st.markdown(f"💭 {course['description']}")
    st.markdown("---")

    ppt_path = course["path"]

    # Load / cache slide images
    if st.session_state.get("current_course_path") != ppt_path:
        with st.spinner("🔄 Loading slides... this may take a moment on first open."):
            try:
                imgs = get_slide_images(ppt_path)
                st.session_state["slide_images"]       = imgs
                st.session_state["current_course_path"] = ppt_path
                st.session_state["slide_index"]         = 0
            except Exception as e:
                st.error(f"❌ Could not render slides: {e}")
                st.markdown("Please download the file to view it locally.")
                with open(ppt_path, "rb") as f:
                    st.download_button("📥 Download PowerPoint", f,
                                       file_name=course["filename"])
                return

    images      = st.session_state.get("slide_images", [])
    total       = len(images)

    if total == 0:
        st.warning("No slides found in this presentation.")
        return

    # Init slide index
    if "slide_index" not in st.session_state:
        st.session_state["slide_index"] = 0

    idx = st.session_state["slide_index"]

    # ── Navigation bar ──
    col_prev, col_counter, col_next = st.columns([1, 2, 1])

    with col_prev:
        if st.button("◀◀  Previous", disabled=(idx == 0),
                     use_container_width=True):
            st.session_state["slide_index"] -= 1
            st.rerun()

    with col_counter:
        st.markdown(
            f'<div class="slide-counter">Slide {idx + 1} / {total}</div>',
            unsafe_allow_html=True
        )

    with col_next:
        if st.button("Next  ▶▶", disabled=(idx == total - 1),
                     use_container_width=True):
            st.session_state["slide_index"] += 1
            st.rerun()

    # ── Progress bar ──
    pct = int((idx + 1) / total * 100)
    st.markdown(
        f'<div class="prog-wrap"><div class="prog-fill" style="width:{pct}%"></div></div>',
        unsafe_allow_html=True
    )

    # ── Slide image ──
    slide_img = images[idx]

    # Convert PIL → bytes for st.image
    buf = io.BytesIO()
    slide_img.save(buf, format="PNG")
    buf.seek(0)

    st.markdown('<div class="slide-frame">', unsafe_allow_html=True)
    st.image(buf, use_container_width=True)
    st.markdown("</div>", unsafe_allow_html=True)

    # ── Keyboard hint + jump ──
    st.markdown("---")
    jump = st.slider("🎯 Jump to slide", min_value=1, max_value=total,
                     value=idx + 1, key="jump_slider")
    if jump - 1 != idx:
        st.session_state["slide_index"] = jump - 1
        st.rerun()

    # ── Download ──
    with open(ppt_path, "rb") as f:
        st.download_button(
            "📥 Download Original PowerPoint", f,
            file_name=course["filename"],
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    st.markdown("</div>", unsafe_allow_html=True)


# ==============================
# MAIN
# ==============================
def main():
    if "viewing_course" not in st.session_state:
        st.session_state["viewing_course"] = None

    st.markdown("""
        <div style="text-align:center; animation:fadeInUp .8s ease-out">
            <h1>🌸 English Teacher's Platform 🌸</h1>
            <p style="color:#c2185b;font-size:18px">✨ Make learning beautiful and fun! ✨</p>
        </div>
        <div style="text-align:center;margin-bottom:20px;font-size:30px">
            📖 📝 🎓 ✏️ 📕
        </div>
    """, unsafe_allow_html=True)

    with st.sidebar:
        st.markdown("""
            <div style="text-align:center;padding:20px 0">
                <h3 style="color:#ff69b4">✨ Welcome! ✨</h3>
                <div style="font-size:30px">👩‍🏫</div>
            </div>
        """, unsafe_allow_html=True)

        mode = st.radio("Choose your role:",
                        ["👩‍🏫 Teacher", "👧 Student"], index=0)
        st.markdown("---")
        st.caption("🌸 Made with love for English teachers 🌸")

    metadata = load_metadata()

    if st.session_state["viewing_course"] is not None:
        display_presentation(st.session_state["viewing_course"])
    elif mode == "👩‍🏫 Teacher":
        teacher_mode(metadata)
    else:
        student_mode(metadata)


# ==============================
# TEACHER MODE
# ==============================
def teacher_mode(metadata):
    st.markdown('<div class="fade-in">', unsafe_allow_html=True)
    col1, col2 = st.columns([2, 1])

    with col1:
        st.subheader("🌸 Upload New Course")
        lc1, lc2 = st.columns(2)
        with lc1:
            level = st.selectbox("📚 Main Level", ["A", "B", "C"])
        with lc2:
            sub_level = st.selectbox("🎯 Sub-level", ["1", "2", "3"])
        full_level = f"{level}{sub_level}"

        title       = st.text_input("📖 Course Title",
                                    placeholder="e.g., Present Simple Tense")
        description = st.text_area("💭 Description",
                                   placeholder="What will students learn?")
        uploaded    = st.file_uploader("📎 Upload PPT/PPTX File",
                                       type=["ppt", "pptx"])

        if st.button("💖 Save Course", use_container_width=True):
            if title and uploaded:
                folder = Path(f"courses/Level_{level}/{full_level}")
                folder.mkdir(parents=True, exist_ok=True)
                save_path = folder / uploaded.name
                with open(save_path, "wb") as f:
                    f.write(uploaded.getbuffer())

                # Pre-render slides cache
                with st.spinner("⚙️ Pre-processing slides..."):
                    get_slide_images(str(save_path))

                key = f"{full_level}_{uploaded.name}"
                metadata[key] = {
                    "title":       title,
                    "description": description or "No description",
                    "level":       full_level,
                    "path":        str(save_path),
                    "filename":    uploaded.name,
                    "upload_date": time.strftime("%Y-%m-%d %H:%M"),
                }
                save_metadata(metadata)
                st.balloons()
                st.success(f"✨ Course '{title}' saved! ✨")
                time.sleep(0.5)
                st.rerun()
            else:
                st.error("💔 Please add a title and file!")

    with col2:
        st.subheader("📊 Quick Stats")
        st.info(f"📚 **Total courses:** {len(metadata)}")
        if metadata:
            cnt = {}
            for c in metadata.values():
                cnt[c["level"]] = cnt.get(c["level"], 0) + 1
            st.write("**Courses per level:**")
            for lv, n in sorted(cnt.items()):
                st.progress(min(n / 10, 1.0), text=f"Level {lv}: {n} courses")

    st.markdown("---")
    st.subheader("📚 Manage Your Courses")

    if metadata:
        fl = st.selectbox("Filter by level:",
                          ["All"] + sorted({c["level"] for c in metadata.values()}))
        for key, course in metadata.items():
            if fl != "All" and course["level"] != fl:
                continue
            with st.container():
                ca, cb, cc = st.columns([3, 1, 1])
                with ca:
                    st.markdown(f"""
                        <div class="course-card">
                            <strong>📘 {course['title']}</strong><br>
                            <small>🎯 Level {course['level']}</small><br>
                            <small>📅 {course['upload_date']}</small><br>
                            <small>💭 {course['description']}</small>
                        </div>""", unsafe_allow_html=True)
                    if st.button("🎬 View & Present", key=f"view_{key}"):
                        for k in ["slide_images", "current_course_path", "slide_index"]:
                            st.session_state.pop(k, None)
                        st.session_state["viewing_course"] = course
                        st.rerun()
                with cb:
                    with open(course["path"], "rb") as f:
                        st.download_button("📥 Download", f,
                                           file_name=course["filename"],
                                           key=f"down_{key}")
                with cc:
                    if st.button("🗑️ Delete", key=f"del_{key}"):
                        if delete_course(key, course["path"]):
                            st.warning(f"💔 '{course['title']}' deleted")
                            time.sleep(0.5)
                            st.rerun()
    else:
        st.info("🌸 No courses yet. Upload your first course above!")

    st.markdown("</div>", unsafe_allow_html=True)


# ==============================
# STUDENT MODE
# ==============================
def student_mode(metadata):
    st.markdown('<div class="fade-in">', unsafe_allow_html=True)
    st.subheader("🎓 Browse Your Courses")

    c1, c2 = st.columns(2)
    with c1:
        level = st.selectbox("📚 Select Main Level", ["A", "B", "C"])
    with c2:
        sub_level = st.selectbox("🎯 Select Sub-level", ["1", "2", "3"])

    full_level       = f"{level}{sub_level}"
    available_courses = {k: v for k, v in metadata.items()
                         if v["level"] == full_level}

    if available_courses:
        st.success(f"✨ Found {len(available_courses)} course(s) for Level {full_level} ✨")
        for key, course in available_courses.items():
            with st.expander(f"📖 {course['title']}", expanded=True):
                ca, cb = st.columns([2, 1])
                with ca:
                    st.markdown(f"""
                        <div style="background:#fff0f5;padding:15px;border-radius:15px">
                            <strong>💭 Description:</strong><br>{course['description']}<br><br>
                            <strong>📅 Uploaded:</strong> {course['upload_date']}<br>
                            <strong>🎯 Level:</strong> {course['level']}
                        </div>""", unsafe_allow_html=True)
                    if st.button("🎬 View Course", key=f"view_s_{key}"):
                        for k in ["slide_images", "current_course_path", "slide_index"]:
                            st.session_state.pop(k, None)
                        st.session_state["viewing_course"] = course
                        st.rerun()
                with cb:
                    with open(course["path"], "rb") as f:
                        st.download_button(
                            "📥 Download Course", f,
                            file_name=course["filename"],
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True,
                            key=f"dl_s_{key}"
                        )
                if st.button("💡 Get a tip", key=f"tip_{key}"):
                    import random
                    tips = ["✨ Take notes while watching!",
                            "💕 Practice with a friend!",
                            "⭐ Review key vocabulary after!",
                            "🌸 Ask questions if something is unclear!"]
                    st.info(f"💖 Tip: {random.choice(tips)}")
    else:
        st.warning(f"💔 No courses available for Level {full_level} yet.")
        st.markdown("""
            <div style="text-align:center;padding:40px">
                <div style="font-size:50px">📚✨</div>
                <p style="color:#c2185b">Ask your teacher to upload courses for this level!</p>
            </div>""", unsafe_allow_html=True)

    st.markdown("</div>", unsafe_allow_html=True)


# ==============================
# ENTRY POINT
# ==============================
if __name__ == "__main__":
    init_folders()
    main()
